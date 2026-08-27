# US-003, step zero: author the seven sample files the conversion probes convert.
#
# Authored inside the sandbox rather than committed from a developer machine, because no machine here
# has Word, Excel or PowerPoint, and a sample nobody can regenerate is a fixture that rots. The
# harness downloads what this writes and commits it, so every later run converts byte-identical
# inputs and a tier change means the platform moved rather than the input did.
#
# The known limitation, recorded rather than papered over: these are authored by the same libraries
# that will read them back, so they exercise the conversion path but not the messiness of a real
# document produced by Office itself. A tier confirmed here is a floor, not a ceiling.
import json
import os

MOUNT = "/mnt/data"
STEM = "sample"

HEADINGS = ["Quarterly Summary", "Regional Performance", "Actions"]
PARAGRAPHS = [
    "Revenue rose across every region except the north, where a single delayed contract accounts for the whole shortfall.",
    "Headcount was flat. Attrition remained inside the planning band for the third consecutive quarter.",
    "The actions below are owned by the regional leads and are reviewed at the next operating meeting.",
]
TABLE = [
    ["Region", "Revenue", "Growth"],
    ["North", "1200", "-4"],
    ["South", "2400", "11"],
    ["East", "1810", "6"],
    ["West", "2050", "9"],
]


def main():
    written = {}




    def record(path):
        written[os.path.splitext(path)[1].lstrip(".")] = {
            "name": os.path.basename(path),
            "bytes": os.path.getsize(path),
        }


    # txt and md, which every other format can be checked against for content survival.
    txt_path = os.path.join(MOUNT, STEM + ".txt")
    with open(txt_path, "w", encoding="utf-8") as handle:
        for heading, paragraph in zip(HEADINGS, PARAGRAPHS):
            handle.write(heading + "\n" + paragraph + "\n\n")
        for row in TABLE:
            handle.write("\t".join(row) + "\n")
    record(txt_path)

    md_path = os.path.join(MOUNT, STEM + ".md")
    with open(md_path, "w", encoding="utf-8") as handle:
        handle.write("# " + HEADINGS[0] + "\n\n")
        handle.write(PARAGRAPHS[0] + "\n\n")
        handle.write("## " + HEADINGS[1] + "\n\n")
        handle.write("| " + " | ".join(TABLE[0]) + " |\n")
        handle.write("| " + " | ".join(["---"] * len(TABLE[0])) + " |\n")
        for row in TABLE[1:]:
            handle.write("| " + " | ".join(row) + " |\n")
        handle.write("\n## " + HEADINGS[2] + "\n\n")
        for paragraph in PARAGRAPHS[1:]:
            handle.write("- " + paragraph + "\n")
    record(md_path)

    csv_path = os.path.join(MOUNT, STEM + ".csv")
    with open(csv_path, "w", encoding="utf-8", newline="") as handle:
        import csv as csv_module

        writer = csv_module.writer(handle)
        for row in TABLE:
            writer.writerow(row)
    record(csv_path)

    # docx: headings, body paragraphs and a ruled table, so a structural conversion has structure to lose.
    import docx

    document = docx.Document()
    for index, heading in enumerate(HEADINGS):
        document.add_heading(heading, level=index + 1)
        document.add_paragraph(PARAGRAPHS[index])
    table = document.add_table(rows=1, cols=len(TABLE[0]))
    table.style = "Table Grid"
    for column, name in enumerate(TABLE[0]):
        table.rows[0].cells[column].text = name
    for row in TABLE[1:]:
        cells = table.add_row().cells
        for column, value in enumerate(row):
            cells[column].text = value
    docx_path = os.path.join(MOUNT, STEM + ".docx")
    document.save(docx_path)
    record(docx_path)

    # xlsx: two sheets, so a converter that silently keeps only the first has somewhere to be caught.
    import openpyxl

    workbook = openpyxl.Workbook()
    sheet = workbook.active
    sheet.title = "Regions"
    for row in TABLE:
        sheet.append(row)
    notes = workbook.create_sheet("Notes")
    for index, paragraph in enumerate(PARAGRAPHS):
        notes.cell(row=index + 1, column=1, value=paragraph)
    xlsx_path = os.path.join(MOUNT, STEM + ".xlsx")
    workbook.save(xlsx_path)
    record(xlsx_path)

    # pptx: three slides with titles and bodies.
    from pptx import Presentation

    presentation = Presentation()
    for heading, paragraph in zip(HEADINGS, PARAGRAPHS):
        slide = presentation.slides.add_slide(presentation.slide_layouts[1])
        slide.shapes.title.text = heading
        slide.placeholders[1].text = paragraph
    pptx_path = os.path.join(MOUNT, STEM + ".pptx")
    presentation.save(pptx_path)
    record(pptx_path)

    # pdf: text-based and carrying a ruled table, which is what pdfplumber's table extraction needs to
    # have any chance. A scanned PDF would be a different probe and a different claim.
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.units import inch
    from reportlab.pdfgen import canvas as pdf_canvas

    pdf_path = os.path.join(MOUNT, STEM + ".pdf")
    page = pdf_canvas.Canvas(pdf_path, pagesize=LETTER)
    width, height = LETTER
    cursor = height - inch
    for heading, paragraph in zip(HEADINGS, PARAGRAPHS):
        page.setFont("Helvetica-Bold", 14)
        page.drawString(inch, cursor, heading)
        cursor -= 0.3 * inch
        page.setFont("Helvetica", 10)
        page.drawString(inch, cursor, paragraph[:95])
        cursor -= 0.5 * inch
    top = cursor
    for row_index, row in enumerate(TABLE):
        y = top - row_index * 0.3 * inch
        for column_index, value in enumerate(row):
            x = inch + column_index * 1.6 * inch
            page.setFont("Helvetica-Bold" if row_index == 0 else "Helvetica", 10)
            page.drawString(x + 4, y + 6, value)
            page.rect(x, y, 1.6 * inch, 0.3 * inch)
    page.showPage()
    page.save()
    record(pdf_path)

    result = {"samples": written}

    with open(os.path.join(MOUNT, "samples.json"), "w", encoding="utf-8") as handle:
        json.dump(result, handle, indent=2)

    print("RESULT:" + json.dumps(result))


main()
