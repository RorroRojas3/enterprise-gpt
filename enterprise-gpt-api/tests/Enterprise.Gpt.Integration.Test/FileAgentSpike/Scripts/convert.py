# US-003: attempt every conversion the matrix proposes for one source format, and record what each
# attempt actually produced.
#
# Three rules this script exists to enforce:
#
#  1. A proposed "refused" pair is attempted anyway. A refusal inherited from documentation about a
#     different environment is not evidence; a refusal recorded after the attempt is.
#  2. Nothing here decides a fidelity tier. It records the engine that ran, whether an output file
#     appeared, and how much of the source survived - the harness verifies the bytes host-side and
#     assigns the tier. A script that graded its own output would grade it optimistically.
#  3. Where a headless office suite is on PATH, the Office-to-PDF path runs through it as well, so the
#     structural and faithful renderings can be compared rather than assumed.
import html as html_module
import json
import os
import re
import subprocess
import traceback

SOURCE_FORMAT = "__SOURCE_FORMAT__"
TARGETS = __TARGETS__
OFFICE_BINARY = "__OFFICE_BINARY__"

MOUNT = "/mnt/data"
OUT_PREFIX = "converted-"

# Structure a converter had to drop because the library that would have carried it is absent. The
# caller grades a cell down on this: an output that parses is not the same as an output that kept
# what the source expressed.
DEGRADED = []


SAMPLE_STEM = "sample"


def find_source(extension):
    # The authored sample by name first. Earlier probes leave their own files in this container - an
    # uploaded marker, a written artifact - and any of them would sort ahead of the sample and be
    # converted instead, which would quietly grade the matrix against the wrong input.
    preferred = os.path.join(MOUNT, SAMPLE_STEM + "." + extension)
    if os.path.exists(preferred):
        return preferred
    candidates = [
        name
        for name in sorted(os.listdir(MOUNT))
        if name.lower().endswith("." + extension) and not name.startswith(OUT_PREFIX)
    ]
    if not candidates:
        raise FileNotFoundError("no ." + extension + " file is mounted under " + MOUNT)
    return os.path.join(MOUNT, candidates[0])


def verify_output(target, path):
    # The second, deterministic, model-free pass FR-25 describes, run here so a tier rests on a file
    # something could actually open rather than on the absence of a traceback. The Office formats are
    # checked twice on purpose: once as an OPC package with the stdlib, which is independent of the
    # library that wrote them, and once through the reader an application would really use.
    info = {"openable": False, "detail": None}
    try:
        if target in ("docx", "xlsx", "pptx"):
            import zipfile

            required = {
                "docx": "word/document.xml",
                "xlsx": "xl/workbook.xml",
                "pptx": "ppt/presentation.xml",
            }[target]
            with zipfile.ZipFile(path) as package:
                names = package.namelist()
            if required not in names:
                info["detail"] = "package carries no " + required
                return info
            if target == "docx":
                import docx

                document = docx.Document(path)
                info["detail"] = "paragraphs=" + str(len(document.paragraphs)) + " tables=" + str(len(document.tables))
            elif target == "xlsx":
                import openpyxl

                workbook = openpyxl.load_workbook(path)
                info["detail"] = "sheets=" + ",".join(workbook.sheetnames)
            else:
                from pptx import Presentation

                info["detail"] = "slides=" + str(len(Presentation(path).slides))
        elif target == "pdf":
            import fitz

            with fitz.open(path) as document:
                pages = document.page_count
                characters = sum(len(page.get_text()) for page in document)
            info["detail"] = "pages=" + str(pages) + " characters=" + str(characters)
        elif target == "csv":
            import csv

            with open(path, newline="", encoding="utf-8") as handle:
                rows = list(csv.reader(handle))
            info["detail"] = "rows=" + str(len(rows))
        else:
            with open(path, encoding="utf-8") as handle:
                text = handle.read()
            info["detail"] = "characters=" + str(len(text))
        info["openable"] = True
    except Exception as error:  # noqa: BLE001 - an unopenable output is the finding
        info["detail"] = type(error).__name__ + ": " + str(error)[:200]
    return info


def out_path(target):
    return os.path.join(MOUNT, OUT_PREFIX + SOURCE_FORMAT + "-to-" + target + "." + target)


def to_markdown_table(frame):
    # Hand-rolled rather than DataFrame.to_markdown, which needs tabulate - absent from this image.
    header = [str(name) for name in frame.columns]
    rows = [[str(value) for value in record] for _, record in frame.iterrows()]
    lines = ["| " + " | ".join(header) + " |", "| " + " | ".join(["---"] * len(header)) + " |"]
    lines.extend("| " + " | ".join(row) + " |" for row in rows)
    return "\n".join(lines)


def escape(text):
    # An unescaped ampersand or angle bracket in the source would break the HTML the renderer parses,
    # which would be recorded as the conversion failing rather than as this script being careless.
    return html_module.escape(str(text))


def write_text(path, text):
    with open(path, "w", encoding="utf-8") as handle:
        handle.write(text)


def html_blocks(html):
    fragments = re.split(r"</(?:p|h1|h2|h3|li|tr|section|div)>", html)
    blocks = []
    for fragment in fragments:
        text = html_module.unescape(re.sub(r"<[^>]+>", " ", fragment)).strip()
        text = re.sub(r"\s+", " ", text)
        if text:
            blocks.append(text)
    return blocks


def reportlab_pdf(html, path):
    from reportlab.lib.pagesizes import LETTER
    from reportlab.lib.styles import getSampleStyleSheet
    from reportlab.platypus import Paragraph, SimpleDocTemplate, Spacer

    styles = getSampleStyleSheet()
    flow = []
    for block in html_blocks(html):
        flow.append(Paragraph(html_module.escape(block), styles["BodyText"]))
        flow.append(Spacer(1, 6))
    SimpleDocTemplate(path, pagesize=LETTER).build(flow or [Paragraph("(empty)", styles["BodyText"])])


def html_to_pdf(html, path):
    # weasyprint is installed here but raises against the image's own Pillow, so the structural path
    # falls through to reportlab rather than reporting every Office-to-PDF pair unachievable.
    try:
        from weasyprint import HTML

        HTML(string=html).write_pdf(path)
        return "weasyprint"
    except Exception as error:  # noqa: BLE001 - which renderer worked is the finding
        reportlab_pdf(html, path)
        return "reportlab (weasyprint failed: " + type(error).__name__ + ")"


def office_to_pdf(source, path):
    if not OFFICE_BINARY:
        raise RuntimeError("no soffice/libreoffice binary is on PATH")
    directory = os.path.dirname(path)
    subprocess.run(
        [OFFICE_BINARY, "--headless", "--convert-to", "pdf", "--outdir", directory, source],
        capture_output=True,
        check=True,
        timeout=300,
    )
    produced = os.path.join(directory, os.path.splitext(os.path.basename(source))[0] + ".pdf")
    if produced != path:
        os.replace(produced, path)
    return os.path.basename(OFFICE_BINARY)


def read_docx(source):
    import docx

    document = docx.Document(source)
    blocks = [paragraph.text for paragraph in document.paragraphs if paragraph.text.strip()]
    tables = [
        ["\t".join(cell.text for cell in row.cells) for row in table.rows] for table in document.tables
    ]
    return blocks, tables


def read_pptx(source):
    from pptx import Presentation

    slides = []
    for slide in Presentation(source).slides:
        lines = []
        for shape in slide.shapes:
            if shape.has_text_frame:
                lines.extend(run.text for run in shape.text_frame.paragraphs[0].runs)
                lines.extend(p.text for p in shape.text_frame.paragraphs if p.text.strip())
        slides.append([line for line in dict.fromkeys(lines) if line.strip()])
    return slides


def read_pdf_text(source):
    import fitz

    with fitz.open(source) as document:
        return [page.get_text() for page in document]


def read_pdf_tables(source):
    import pdfplumber

    tables = []
    with pdfplumber.open(source) as document:
        for page in document.pages:
            tables.extend(page.extract_tables() or [])
    return tables


def docx_to(target, source, path):
    blocks, tables = read_docx(source)
    if target in ("md", "txt"):
        body = "\n\n".join(blocks)
        for table in tables:
            body += "\n\n" + "\n".join(table)
        write_text(path, body)
        return "python-docx", {"paragraphs": len(blocks), "tables": len(tables)}
    if target == "pdf":
        if OFFICE_BINARY:
            return office_to_pdf(source, path), {"paragraphs": len(blocks)}
        html = "".join("<p>" + escape(block) + "</p>" for block in blocks)
        return html_to_pdf("<html><body>" + html + "</body></html>", path), {"paragraphs": len(blocks)}
    raise NotImplementedError(SOURCE_FORMAT + " to " + target)


def xlsx_to(target, source, path):
    import pandas

    sheets = pandas.read_excel(source, sheet_name=None)
    first = next(iter(sheets.values()))
    metrics = {"sheets": len(sheets), "rows": int(first.shape[0]), "columns": int(first.shape[1])}
    if target == "csv":
        first.to_csv(path, index=False)
        return "pandas", metrics
    if target == "md":
        write_text(path, to_markdown_table(first))
        return "pandas", metrics
    if target == "txt":
        write_text(path, first.to_string(index=False))
        return "pandas", metrics
    if target == "pdf":
        if OFFICE_BINARY:
            return office_to_pdf(source, path), metrics
        return html_to_pdf("<html><body>" + first.to_html(index=False) + "</body></html>", path), metrics
    raise NotImplementedError(SOURCE_FORMAT + " to " + target)


def pptx_to(target, source, path):
    slides = read_pptx(source)
    metrics = {"slides": len(slides)}
    if target in ("md", "txt"):
        separator = "\n\n---\n\n" if target == "md" else "\n\n"
        write_text(path, separator.join("\n".join(slide) for slide in slides))
        return "python-pptx", metrics
    if target == "pdf":
        if OFFICE_BINARY:
            return office_to_pdf(source, path), metrics
        html = "".join(
            "<section><h2>Slide " + str(index + 1) + "</h2>"
            + "".join("<p>" + escape(line) + "</p>" for line in slide)
            + "</section>"
            for index, slide in enumerate(slides)
        )
        return html_to_pdf("<html><body>" + html + "</body></html>", path), metrics
    raise NotImplementedError(SOURCE_FORMAT + " to " + target)


def pdf_to(target, source, path):
    if target in ("md", "txt"):
        pages = read_pdf_text(source)
        write_text(path, "\n\n".join(pages))
        return "PyMuPDF", {"pages": len(pages), "characters": sum(len(page) for page in pages)}
    if target == "docx":
        import docx

        pages = read_pdf_text(source)
        document = docx.Document()
        for page in pages:
            for line in page.splitlines():
                if line.strip():
                    document.add_paragraph(line)
        document.save(path)
        return "PyMuPDF + python-docx", {"pages": len(pages)}
    if target in ("csv", "xlsx"):
        tables = read_pdf_tables(source)
        if target == "csv":
            import csv

            with open(path, "w", encoding="utf-8", newline="") as handle:
                writer = csv.writer(handle)
                for row in (tables[0] if tables else []):
                    writer.writerow(row)
        else:
            import openpyxl

            workbook = openpyxl.Workbook()
            sheet = workbook.active
            for row in (tables[0] if tables else []):
                sheet.append(row)
            workbook.save(path)
        return "pdfplumber", {"tablesFound": len(tables)}
    if target == "pptx":
        # Proposed "refused". Attempted anyway, and what it produces is the evidence for the refusal:
        # one rasterised page per slide with no editable text on it.
        import fitz
        from pptx import Presentation
        from pptx.util import Inches

        presentation = Presentation()
        with fitz.open(source) as document:
            for index, page in enumerate(document):
                image = os.path.join(MOUNT, "page-" + str(index) + ".png")
                page.get_pixmap(dpi=96).save(image)
                slide = presentation.slides.add_slide(presentation.slide_layouts[6])
                slide.shapes.add_picture(image, Inches(0), Inches(0), width=presentation.slide_width)
            pages = document.page_count
        presentation.save(path)
        return "PyMuPDF + python-pptx", {"pages": pages, "editableTextShapes": 0}
    raise NotImplementedError(SOURCE_FORMAT + " to " + target)


def csv_to(target, source, path):
    import pandas

    frame = pandas.read_csv(source)
    metrics = {"rows": int(frame.shape[0]), "columns": int(frame.shape[1])}
    if target == "xlsx":
        frame.to_excel(path, index=False)
        return "pandas + openpyxl", metrics
    if target == "md":
        write_text(path, to_markdown_table(frame))
        return "pandas", metrics
    if target == "txt":
        write_text(path, frame.to_string(index=False))
        return "pandas", metrics
    if target == "docx":
        import docx

        document = docx.Document()
        table = document.add_table(rows=1, cols=frame.shape[1])
        for column, name in enumerate(frame.columns):
            table.rows[0].cells[column].text = str(name)
        for _, record in frame.iterrows():
            cells = table.add_row().cells
            for column, value in enumerate(record):
                cells[column].text = str(value)
        document.save(path)
        return "pandas + python-docx", metrics
    if target == "pdf":
        return html_to_pdf("<html><body>" + frame.to_html(index=False) + "</body></html>", path), metrics
    raise NotImplementedError(SOURCE_FORMAT + " to " + target)


def markup_to(target, source, path):
    with open(source, "r", encoding="utf-8") as handle:
        text = handle.read()

    metrics = {"characters": len(text), "lines": len(text.splitlines())}

    if target in ("md", "txt"):
        write_text(path, text)
        return "stdlib", metrics
    if target == "docx":
        import docx

        document = docx.Document()
        for line in text.splitlines():
            if line.strip():
                document.add_paragraph(line)
        document.save(path)
        return "stdlib + python-docx", metrics
    if target == "pdf":
        # Imported here, not at the top: markdown is not in this image, and an import that sinks the
        # plain-text targets would report them as unachievable when only the HTML path needs it.
        engine = "stdlib"
        html = None
        if SOURCE_FORMAT == "md":
            try:
                import markdown

                html = markdown.markdown(text, extensions=["tables"])
                engine = "markdown"
            except ImportError:
                DEGRADED.append(
                    "markdown is absent, so headings, lists and tables were flattened to paragraphs"
                )
                html = None
        if html is None:
            html = "".join("<p>" + escape(line) + "</p>" for line in text.splitlines() if line.strip())
        return engine + " + " + html_to_pdf("<html><body>" + html + "</body></html>", path), metrics
    raise NotImplementedError(SOURCE_FORMAT + " to " + target)


CONVERTERS = {
    "docx": docx_to,
    "xlsx": xlsx_to,
    "pptx": pptx_to,
    "pdf": pdf_to,
    "csv": csv_to,
    "md": markup_to,
    "txt": markup_to,
}

# Resolved inside a try rather than at module scope. A container recycled between the run that
# authored the samples and this one leaves nothing to convert, and dying here would print no RESULT
# line at all - so the caller would see "the probe did not complete" after several minutes of billed
# sandbox time, with no way to tell that apart from a real failure. Reported instead, so the caller
# can re-author the corpus and retry once.
source_path = None
source_error = None

try:
    source_path = find_source(SOURCE_FORMAT)
except Exception as error:  # noqa: BLE001 - a missing source is a reportable state, not a crash
    source_error = type(error).__name__ + ": " + str(error)[:300]

convert = CONVERTERS[SOURCE_FORMAT]
attempts = []

for target in TARGETS if source_path is not None else []:
    path = out_path(target)
    attempt = {"from": SOURCE_FORMAT, "to": target, "output": os.path.basename(path)}
    DEGRADED.clear()
    try:
        engine, metrics = convert(target, source_path, path)
        attempt["produced"] = os.path.exists(path)
        attempt["bytes"] = os.path.getsize(path) if attempt["produced"] else 0
        attempt["engine"] = engine
        attempt["metrics"] = metrics
        attempt["error"] = None
    except Exception as error:  # noqa: BLE001 - the failure IS the evidence for a refusal
        attempt["produced"] = False
        attempt["bytes"] = 0
        attempt["engine"] = None
        attempt["metrics"] = {}
        attempt["error"] = type(error).__name__ + ": " + str(error)[:300]
        attempt["traceback"] = traceback.format_exc()[-800:]
    attempt["degraded"] = list(DEGRADED)
    attempt["verification"] = (
        verify_output(target, path)
        if attempt["produced"]
        else {"openable": False, "detail": "nothing was produced"}
    )
    attempts.append(attempt)

result = {
    "source": SOURCE_FORMAT,
    "sourceFound": source_path is not None,
    "sourceFile": os.path.basename(source_path) if source_path is not None else None,
    "sourceBytes": os.path.getsize(source_path) if source_path is not None else 0,
    "sourceError": source_error,
    "officeBinary": OFFICE_BINARY or None,
    "attempts": attempts,
}

with open(os.path.join(MOUNT, "conversion-" + SOURCE_FORMAT + ".json"), "w", encoding="utf-8") as handle:
    json.dump(result, handle, indent=2)

print("RESULT:" + json.dumps(result))
